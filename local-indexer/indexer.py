"""
Knowledge Wiki - Local Indexer v2.0
====================================
Google Drive 동기화 폴더의 문서를 파싱하여 Knowledge Wiki에 업로드합니다.
v2.0: 자동 메타데이터 태깅 (키워드 추출, 주제분류, 기관 인식, 문서단계 판별)

지원 파일:
  - PPTX (슬라이드 단위)
  - PDF (페이지 단위)
  - XLSX (시트+행 단위)
  - CSV (행 단위)
  - ipynb (셀 단위)
  - DOCX (문단 단위)

사용법:
  1. pip install -r requirements.txt
  2. config.py에서 DRIVE_ROOT, WIKI_API_URL 설정
  3. python indexer.py
"""

import os
import sys
import json
import hashlib
import sqlite3
import re
import time
import traceback
from pathlib import Path
from datetime import datetime
from collections import Counter

import requests

from config import DRIVE_ROOT, WIKI_API_URL, BATCH_SIZE, SUPPORTED_EXTENSIONS


# =============================================
# Auto-Tagging Configuration
# =============================================

# 주제분류 키워드 매핑
CATEGORY_KEYWORDS = {
    'AI': ['AI', '인공지능', 'LLM', '생성형', '머신러닝', '딥러닝', 'GPT', 'RAG', '멀티모달',
           '챗봇', 'NLP', '자연어', '강화학습', 'XAI', '연합학습', '에이전트', 'transformer'],
    '데이터': ['데이터', '빅데이터', 'data', '공공데이터', '데이터셋', '오픈데이터', '데이터 개방',
              '데이터 품질', '데이터 표준', '데이터 유통', '마이데이터', '데이터 가치'],
    '거버넌스': ['거버넌스', 'CDO', '데이터 관리', '메타데이터', '마스터데이터', '표준화',
               '품질관리', '성숙도', '거버넌스 체계', 'EA', '아키텍처'],
    '인프라': ['인프라', '서버', '클라우드', '네트워크', '스토리지', 'AWS', 'GCP', 'Azure',
             '온프레미스', 'CDN', 'API', '마이크로서비스', '컨테이너', 'Docker', 'K8s'],
    '보안': ['보안', '개인정보', '비식별', '암호화', '방화벽', 'IPS', 'WAF', '접근제어',
            'ISMS', '정보보호', 'GDPR', '프라이버시', '인증', 'SSO'],
    '전략': ['전략', 'ISP', 'ISMP', '로드맵', '비전', 'KPI', '목표모델', 'To-Be', 'As-Is',
            '이행계획', '마스터플랜', '중장기', '디지털 전환', 'DX'],
    '사업관리': ['RFP', '제안', '과업', '예산', '일정', '평가', '검수', '착수', '보고', '산출물',
               'WBS', 'PMO', '리스크', '품질보증'],
}

# 문서단계 판별 키워드
DOC_STAGE_PATTERNS = {
    'RFP': ['RFP', '제안요청서', '제안요청', '과업지시서', '입찰공고'],
    '제안서': ['제안서', '기술제안', '사업제안', '수행계획'],
    '착수보고': ['착수보고', '착수', '킥오프', 'kickoff'],
    '중간보고': ['중간보고', '중간점검', '진행보고'],
    '최종보고': ['최종보고', '완료보고', '결과보고'],
    '산출물': ['산출물', '결과물', '분석보고서', '현황분석'],
    '조사자료': ['조사', '설문', '현황', '목록', '리스트', 'survey'],
    '분석코드': ['분석코드', 'analysis', 'notebook', '.ipynb'],
}

# 기관 사전
ORG_DICT = {
    'NIA': ['NIA', '한국지능정보사회진흥원', '정보사회진흥원'],
    '과학기술정보통신부': ['과기정통부', '과학기술정보통신부', '과학기술부'],
    '행정안전부': ['행안부', '행정안전부'],
    '보건복지부': ['보건복지부', '복지부'],
    '국토교통부': ['국토교통부', '국토부'],
    '환경부': ['환경부'],
    '한국정보화진흥원': ['한국정보화진흥원', '정보화진흥원'],
    'NIPA': ['NIPA', '정보통신산업진흥원'],
    'KEIT': ['KEIT', '한국산업기술평가관리원'],
    'KDATA': ['KDATA', '한국데이터산업진흥원'],
}

# 불용어
STOP_WORDS = set([
    '있다', '없다', '것이다', '한다', '등', '및', '의', '을', '를', '이', '가', '에',
    '위한', '대한', '통한', '따른', '필요', '관련', '경우', '사항', '내용', '현황',
    '수행', '추진', '운영', '구축', '개발', '활용', '지원', '제공', '확대', '강화',
    '개선', '분석', '설계', '관리', '기반', '시스템', '서비스', '정보', '기관', '사업',
    'the', 'and', 'for', 'with', 'from', 'this', 'that', 'are', 'was', 'were',
])


# =============================================
# Auto-Tagging Functions
# =============================================

def extract_tags(text, max_tags=8):
    """TF 기반 키워드 추출 (단순 빈도)"""
    if not text:
        return []
    
    # 한글+영문 토큰 추출 (2글자 이상)
    tokens = re.findall(r'[가-힣]{2,}|[A-Za-z]{2,}', text)
    tokens = [t for t in tokens if t.lower() not in STOP_WORDS and len(t) >= 2]
    
    # 빈도 카운트
    counter = Counter(tokens)
    
    # 상위 키워드 반환
    return [word for word, _ in counter.most_common(max_tags)]


def classify_category(text, filepath=''):
    """텍스트 기반 주제분류"""
    if not text:
        return '', ''
    
    combined = text + ' ' + filepath
    scores = {}
    
    for cat, keywords in CATEGORY_KEYWORDS.items():
        score = 0
        for kw in keywords:
            count = combined.lower().count(kw.lower())
            score += count
        if score > 0:
            scores[cat] = score
    
    if not scores:
        return '', ''
    
    # 최고 점수 카테고리
    sorted_cats = sorted(scores.items(), key=lambda x: x[1], reverse=True)
    main_cat = sorted_cats[0][0]
    sub_cat = sorted_cats[1][0] if len(sorted_cats) > 1 and sorted_cats[1][1] > 2 else ''
    
    return main_cat, sub_cat


def detect_doc_stage(filepath, text=''):
    """파일경로 + 텍스트로 문서단계 판별"""
    combined = filepath + ' ' + (text[:500] if text else '')
    
    for stage, patterns in DOC_STAGE_PATTERNS.items():
        for pat in patterns:
            if pat.lower() in combined.lower():
                return stage
    
    # 폴더명 기반 추정
    path_parts = filepath.lower()
    folder_stage_map = {
        '01.': 'RFP', '01 ': 'RFP',
        '02.': '산출물', '02 ': '산출물',
        '03.': '최종보고', '03 ': '최종보고',
        '04.': '산출물', '04 ': '산출물',
        '05.': '조사자료', '05 ': '조사자료',
        '06.': '분석코드', '06 ': '분석코드',
    }
    for prefix, stage in folder_stage_map.items():
        if prefix in path_parts:
            return stage
    
    return ''


def detect_org(text, filepath=''):
    """발주기관 인식"""
    combined = text + ' ' + filepath
    
    for org_name, patterns in ORG_DICT.items():
        for pat in patterns:
            if pat in combined:
                return org_name
    
    return ''


def extract_author(text, filepath=''):
    """작성자/수행기관 추출 (표지 텍스트 패턴)"""
    # 일반적인 패턴: "수행기관: XXX", "작성: XXX컨설팅"
    patterns = [
        r'수행기관[:\s]+([가-힣A-Za-z]+(?:컨설팅|연구원|연구소|주식회사|㈜))',
        r'작성[:\s]+([가-힣A-Za-z]+)',
        r'(?:㈜|주\))\s*([가-힣]+)',
    ]
    
    for pat in patterns:
        m = re.search(pat, text[:1000])
        if m:
            return m.group(1).strip()
    
    return ''


def extract_year(text, filepath=''):
    """사업연도 추출"""
    combined = filepath + ' ' + (text[:500] if text else '')
    
    # 2020~2029년도 패턴
    years = re.findall(r'20[2-3]\d', combined)
    if years:
        return max(years)  # 가장 최근 연도
    
    return ''


def calc_importance(text, doc_stage='', filepath=''):
    """문서 중요도 점수 (0-100)"""
    score = 50  # 기본
    
    # 문서 단계 가중치
    stage_weight = {
        '최종보고': 20, '제안서': 15, '산출물': 10,
        'RFP': 8, '중간보고': 5, '착수보고': 3,
        '조사자료': 0, '분석코드': -5,
    }
    score += stage_weight.get(doc_stage, 0)
    
    # 텍스트 길이 가중치 (내용이 풍부할수록)
    text_len = len(text) if text else 0
    if text_len > 500:
        score += 10
    elif text_len > 200:
        score += 5
    
    # 핵심 키워드 포함 시 가중치
    important_keywords = ['전략', '비전', '목표', 'KPI', '로드맵', '핵심', '결론', '권고']
    for kw in important_keywords:
        if kw in (text or ''):
            score += 3
    
    return min(max(score, 10), 100)  # 10~100 범위


def auto_tag_chunk(chunk, filepath):
    """청크에 자동 메타데이터 태깅 적용"""
    text = chunk.get('text', '')
    
    # Tags
    chunk['tags'] = extract_tags(text)
    
    # Category
    cat, sub_cat = classify_category(text, filepath)
    chunk['category'] = cat
    chunk['sub_category'] = sub_cat
    
    # Doc Stage
    chunk['doc_stage'] = detect_doc_stage(filepath, text)
    
    # Organization
    chunk['org'] = detect_org(text, filepath)
    
    # Author
    chunk['author'] = extract_author(text, filepath)
    
    # Year
    chunk['doc_year'] = extract_year(text, filepath)
    
    # Importance
    chunk['importance'] = calc_importance(text, chunk['doc_stage'], filepath)
    
    # Summary (간단한 첫 줄 요약 - 추후 AI 대체)
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    chunk['summary'] = lines[0][:150] if lines else ''
    
    return chunk


# =============================================
# Database (tracking indexed files)
# =============================================
DB_PATH = os.path.join(os.path.dirname(__file__), 'indexer_state.db')

def init_db():
    conn = sqlite3.connect(DB_PATH)
    conn.execute('''
        CREATE TABLE IF NOT EXISTS indexed_files (
            file_path TEXT PRIMARY KEY,
            mtime REAL,
            hash TEXT,
            chunk_count INTEGER DEFAULT 0,
            last_indexed TEXT,
            status TEXT DEFAULT 'ok'
        )
    ''')
    conn.commit()
    return conn

def get_file_hash(filepath, block_size=65536):
    size = os.path.getsize(filepath)
    hasher = hashlib.sha256()
    hasher.update(str(size).encode())
    with open(filepath, 'rb') as f:
        buf = f.read(block_size)
        hasher.update(buf)
        if size > block_size * 2:
            f.seek(-block_size, 2)
            buf = f.read(block_size)
            hasher.update(buf)
    return hasher.hexdigest()[:16]

def needs_indexing(conn, filepath, mtime):
    row = conn.execute(
        'SELECT mtime, hash FROM indexed_files WHERE file_path = ?', (filepath,)
    ).fetchone()
    if row is None:
        return True
    if abs(row[0] - mtime) > 1:
        return True
    return False

def mark_indexed(conn, filepath, mtime, file_hash, chunk_count, status='ok'):
    conn.execute('''
        INSERT OR REPLACE INTO indexed_files (file_path, mtime, hash, chunk_count, last_indexed, status)
        VALUES (?, ?, ?, ?, ?, ?)
    ''', (filepath, mtime, file_hash, chunk_count, datetime.now().isoformat(), status))
    conn.commit()


# =============================================
# Parsers (unchanged from v1)
# =============================================

def parse_pptx(filepath):
    from pptx import Presentation
    chunks = []
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(' | '.join(row_texts))
        full_text = '\n'.join(texts)
        if full_text.strip():
            chunks.append({
                'location_type': 'slide',
                'location_value': str(i),
                'location_detail': f'Slide {i}',
                'text': full_text
            })
    return chunks

def parse_pdf(filepath):
    import fitz
    chunks = []
    doc = fitz.open(filepath)
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            chunks.append({
                'location_type': 'page',
                'location_value': str(i),
                'location_detail': f'Page {i}',
                'text': text
            })
    doc.close()
    return chunks

def parse_xlsx(filepath):
    from openpyxl import load_workbook
    chunks = []
    wb = load_workbook(filepath, read_only=True, data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        row_num = 0
        for row in ws.iter_rows(values_only=True):
            row_num += 1
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                text = ' | '.join(cells)
                chunks.append({
                    'location_type': 'sheet',
                    'location_value': sheet_name,
                    'location_detail': f'Sheet:{sheet_name} Row:{row_num}',
                    'text': text
                })
    wb.close()
    return chunks

def parse_csv(filepath):
    import csv
    chunks = []
    for enc in ['utf-8', 'cp949', 'euc-kr', 'latin-1']:
        try:
            with open(filepath, 'r', encoding=enc) as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader, 1):
                    cells = [c.strip() for c in row if c.strip()]
                    if cells:
                        text = ','.join(cells)
                        chunks.append({
                            'location_type': 'row',
                            'location_value': str(i),
                            'location_detail': f'Row {i}',
                            'text': text
                        })
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
    return chunks

def parse_ipynb(filepath):
    chunks = []
    with open(filepath, 'r', encoding='utf-8') as f:
        nb = json.load(f)
    for i, cell in enumerate(nb.get('cells', []), 1):
        cell_type = cell.get('cell_type', 'code')
        source = ''.join(cell.get('source', []))
        if source.strip():
            chunks.append({
                'location_type': 'cell',
                'location_value': str(i),
                'location_detail': f'Cell {i} ({cell_type})',
                'text': source
            })
    return chunks

def parse_docx(filepath):
    from docx import Document
    chunks = []
    doc = Document(filepath)
    current_text = []
    para_start = 1
    for i, para in enumerate(doc.paragraphs, 1):
        text = para.text.strip()
        if text:
            current_text.append(text)
        if len(current_text) >= 5 or (i == len(doc.paragraphs) and current_text):
            chunks.append({
                'location_type': 'page',
                'location_value': str(para_start),
                'location_detail': f'Paragraphs {para_start}-{i}',
                'text': '\n'.join(current_text)
            })
            current_text = []
            para_start = i + 1
    return chunks


PARSERS = {
    '.pptx': parse_pptx, '.pdf': parse_pdf, '.xlsx': parse_xlsx,
    '.csv': parse_csv, '.ipynb': parse_ipynb, '.docx': parse_docx,
}


# =============================================
# Scanner & Upload
# =============================================

def scan_files(root_dir):
    files = []
    root = Path(root_dir)
    for ext in SUPPORTED_EXTENSIONS:
        for f in root.rglob(f'*{ext}'):
            if '~$' in f.name or '.git' in str(f) or 'node_modules' in str(f):
                continue
            files.append(f)
    return files

def get_project_path(filepath, root_dir):
    rel = os.path.relpath(filepath, root_dir)
    parts = Path(rel).parts
    return parts[0] if len(parts) > 1 else ''

def upload_chunks(chunks, api_url):
    url = f'{api_url}/api/chunks'
    for i in range(0, len(chunks), BATCH_SIZE):
        batch = chunks[i:i + BATCH_SIZE]
        try:
            resp = requests.post(url, json={'chunks': batch}, timeout=30)
            data = resp.json()
            if data.get('errors'):
                print(f"  [WARN] Upload errors: {data['errors']}")
            return data.get('inserted', 0)
        except Exception as e:
            print(f"  [ERROR] Upload failed: {e}")
            return 0
    return 0


# =============================================
# Main
# =============================================

def main():
    print("=" * 60)
    print("  Knowledge Wiki - Local Indexer v2.0")
    print("  (with Auto Metadata Tagging)")
    print("=" * 60)
    print(f"  Drive Root: {DRIVE_ROOT}")
    print(f"  Wiki API:   {WIKI_API_URL}")
    print()

    if not os.path.isdir(DRIVE_ROOT):
        print(f"[ERROR] Drive root not found: {DRIVE_ROOT}")
        print("  Please check config.py and set the correct path.")
        sys.exit(1)

    conn = init_db()

    # Scan
    print("[1/4] Scanning files...")
    files = scan_files(DRIVE_ROOT)
    print(f"  Found {len(files)} supported files")

    # Check changes
    print("[2/4] Checking for changes...")
    to_index = []
    for f in files:
        mtime = os.path.getmtime(f)
        if needs_indexing(conn, str(f), mtime):
            to_index.append(f)
    print(f"  {len(to_index)} files need (re)indexing")
    
    if not to_index:
        print("\n[DONE] Everything is up to date!")
        conn.close()
        return

    # Parse + Auto-tag + Upload
    print(f"[3/4] Parsing & auto-tagging {len(to_index)} files...")
    total_chunks = 0
    errors = 0
    tag_stats = Counter()

    for i, filepath in enumerate(to_index, 1):
        ext = filepath.suffix.lower()
        parser = PARSERS.get(ext)
        if not parser:
            continue

        rel_path = os.path.relpath(filepath, DRIVE_ROOT)
        project = get_project_path(filepath, DRIVE_ROOT)
        doc_title = filepath.stem
        mtime_raw = os.path.getmtime(filepath)
        mtime = datetime.fromtimestamp(mtime_raw).isoformat()
        file_hash = get_file_hash(str(filepath))

        print(f"  [{i}/{len(to_index)}] {rel_path}...", end=' ', flush=True)

        try:
            raw_chunks = parser(str(filepath))
            
            chunks = []
            for c in raw_chunks:
                chunk_id = f"{file_hash}-{c['location_type']}-{c['location_value']}"
                chunk = {
                    'chunk_id': chunk_id,
                    'file_path': rel_path.replace('\\', '/'),
                    'file_type': ext.lstrip('.'),
                    'project_path': project,
                    'doc_title': doc_title,
                    'location_type': c['location_type'],
                    'location_value': c['location_value'],
                    'location_detail': c['location_detail'],
                    'text': c['text'],
                    'mtime': mtime,
                    'hash': file_hash,
                }
                
                # ★ v2.0: Auto-tag each chunk
                chunk = auto_tag_chunk(chunk, rel_path)
                chunks.append(chunk)
                
                # Stats
                if chunk.get('category'):
                    tag_stats[chunk['category']] += 1
            
            if chunks:
                uploaded = upload_chunks(chunks, WIKI_API_URL)
                total_chunks += len(chunks)
                cat_info = chunks[0].get('category', '?')
                stage_info = chunks[0].get('doc_stage', '?')
                print(f"{len(chunks)} chunks [cat:{cat_info} stage:{stage_info}]")
            else:
                print("(empty)")

            mark_indexed(conn, str(filepath), mtime_raw, file_hash, len(chunks))

        except Exception as e:
            print(f"ERROR: {e}")
            mark_indexed(conn, str(filepath), mtime_raw, '', 0, status=f'error: {e}')
            errors += 1

    print(f"\n[4/4] Upload complete!")
    print(f"  Total chunks: {total_chunks}")
    print(f"  Errors: {errors}")
    if tag_stats:
        print(f"  Category distribution: {dict(tag_stats.most_common())}")
    
    conn.close()
    print("\n[DONE]")


if __name__ == '__main__':
    main()
